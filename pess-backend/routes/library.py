from flask import Blueprint, render_template, request, redirect, url_for, flash, session, send_from_directory
import sqlite3, os
from werkzeug.utils import secure_filename
import fitz          # PyMuPDF for PDF
import docx          # python-docx for Word
import pptx          # python-pptx for PowerPoint
import openpyxl      # for Excel

library_bp = Blueprint("library", __name__, url_prefix="/library")
DB_PATH = "pess.db"
UPLOAD_FOLDER = "uploads"

def query_db(query, args=(), one=False):
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.execute(query, args)
    rv = cur.fetchall()
    conn.commit()
    conn.close()
    return (rv[0] if rv else None) if one else rv

# --- Text Extraction Helpers ---
def extract_text_from_pdf(filepath):
    text = ""
    doc = fitz.open(filepath)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(filepath):
    doc = docx.Document(filepath)
    return " ".join([p.text for p in doc.paragraphs])

def extract_text_from_txt(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()

def extract_text_from_pptx(filepath):
    prs = pptx.Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_text_from_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell]) + " "
    return text

def extract_text(filepath, ext):
    if ext == ".pdf": return extract_text_from_pdf(filepath)
    elif ext == ".docx": return extract_text_from_docx(filepath)
    elif ext in [".txt", ".md"]: return extract_text_from_txt(filepath)
    elif ext == ".pptx": return extract_text_from_pptx(filepath)
    elif ext == ".xlsx": return extract_text_from_xlsx(filepath)
    else: return None

# --- AI Helpers ---
def summarize_text(text, max_sentences=3):
    sentences = text.split(".")
    return ". ".join(sentences[:max_sentences]) + "..."

def generate_quiz(text, num_questions=3):
    words = text.split()
    questions = []
    for i in range(min(num_questions, len(words)//10)):
        keyword = words[i*10].capitalize()
        questions.append({
            "q": f"What is the importance of {keyword}?",
            "a": f"{keyword} relates to study material context."
        })
    return questions

def generate_flashcards(text, num_cards=5):
    words = text.split()
    cards = []
    for i in range(min(num_cards, len(words)//15)):
        keyword = words[i*15].capitalize()
        cards.append({
            "front": f"What is {keyword}?",
            "back": f"{keyword} is explained in the study material."
        })
    return cards

# --- Routes ---
@library_bp.route("/")
def library_home():
    files = query_db("SELECT * FROM library")
    links = query_db("SELECT * FROM safe_links")
    return render_template("library.html", files=files, links=links)

@library_bp.route("/upload", methods=["POST"])
def upload_file():
    file = request.files["file"]
    if file:
        filename = secure_filename(file.filename)
        ext = os.path.splitext(filename)[1].lower()
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)

        conn = sqlite3.connect(DB_PATH)
        cur = conn.cursor()
        cur.execute("INSERT INTO library (filename, filepath, extension, uploaded_by) VALUES (?, ?, ?, ?)",
                    (filename, filepath, ext, session.get("username")))
        conn.commit()
        conn.close()
        flash(f"File '{filename}' uploaded successfully!", "success")
    return redirect(url_for("library.library_home"))

@library_bp.route("/delete/<int:file_id>")
def delete_file(file_id):
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("DELETE FROM library WHERE id=?", (file_id,))
    conn.commit()
    conn.close()
    flash("File deleted successfully!", "info")
    return redirect(url_for("library.library_home"))

@library_bp.route("/rename/<int:file_id>", methods=["POST"])
def rename_file(file_id):
    new_name = request.form["new_name"]
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    cur.execute("UPDATE library SET filename=? WHERE id=?", (new_name, file_id))
    conn.commit()
    conn.close()
    flash("File renamed successfully!", "success")
    return redirect(url_for("library.library_home"))

@library_bp.route("/download/<int:file_id>")
def download_file(file_id):
    file = query_db("SELECT * FROM library WHERE id=?", (file_id,), one=True)
    if file:
        return send_from_directory(os.path.dirname(file["filepath"]), os.path.basename(file["filepath"]), as_attachment=True)
    flash("File not found!", "danger")
    return redirect(url_for("library.library_home"))

@library_bp.route("/summarize/<int:file_id>")
def summarize_file(file_id):
    file = query_db("SELECT * FROM library WHERE id=?", (file_id,), one=True)
    if file:
        text = extract_text(file["filepath"], file["extension"])
        if text:
            summary = summarize_text(text)
            return render_template("summary.html", filename=file["filename"], summary=summary)
    flash("Summarization not supported for this file type.", "warning")
    return redirect(url_for("library.library_home"))

@library_bp.route("/quiz/<int:file_id>")
def quiz_from_file(file_id):
    file = query_db("SELECT * FROM library WHERE id=?", (file_id,), one=True)
    if file:
        text = extract_text(file["filepath"], file["extension"])
        if text:
            quiz = generate_quiz(text)
            return render_template("quiz.html", filename=file["filename"], quiz=quiz)
    flash("Quiz not supported for this file type.", "warning")
    return redirect(url_for("library.library_home"))

@library_bp.route("/flashcards/<int:file_id>")
def flashcards_from_file(file_id):
    file = query_db("SELECT * FROM library WHERE id=?", (file_id,), one=True)
    if file:
        text = extract_text(file["filepath"], file["extension"])
        if text:
            cards = generate_flashcards(text)
            return render_template("flashcards.html", filename=file["filename"], cards=cards)
    flash("Flashcards not supported for this file type.", "warning")
    return redirect(url_for("library.library_home"))
